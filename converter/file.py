import os
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from .text import convert_text  # Avtomatik Kiril ↔ Lotin funksiyasi

def convert_docx(file_path, output_folder):
    doc = Document(file_path)
    for para in doc.paragraphs:
        para.text = convert_text(para.text)

    # Shuningdek, jadval hujayralarini ham tarjima qilish
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                cell.text = convert_text(cell.text)

    output_path = os.path.join(output_folder, os.path.basename(file_path))
    doc.save(output_path)
    return output_path

def convert_pptx(file_path, output_folder):
    pres = Presentation(file_path)
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = convert_text(run.text)

    output_path = os.path.join(output_folder, os.path.basename(file_path))
    pres.save(output_path)
    return output_path

def convert_xlsx(file_path, output_folder):
    wb = load_workbook(file_path)
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    cell.value = convert_text(cell.value)
    output_path = os.path.join(output_folder, os.path.basename(file_path))
    wb.save(output_path)
    return output_path

def convert_file(file_path, output_folder):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".docx":
            out = convert_docx(file_path, output_folder)
        elif ext == ".pptx":
            out = convert_pptx(file_path, output_folder)
        elif ext == ".xlsx":
            out = convert_xlsx(file_path, output_folder)
        else:
            return None, None
        return out, os.path.basename(out)
    except Exception as e:
        print(f"[Xatolik] Fayl o‘girishda xato: {e}")
        return None, None
